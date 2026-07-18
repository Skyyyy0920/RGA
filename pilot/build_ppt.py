#!/usr/bin/env python
"""Plain whiteboard slide deck — method explanation, formulas front and center."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

BLACK = RGBColor(0x11, 0x11, 0x11)
GREY  = RGBColor(0x33, 0x33, 0x33)
BLUE  = RGBColor(0x1A, 0x4E, 0x8A)
RED   = RGBColor(0xA1, 0x2A, 0x1F)
GREEN = RGBColor(0x1E, 0x6B, 0x2E)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def _emit(p, latex, font, size, color, bold):
    """Render simple LaTeX-ish text with _sub and ^sup into runs."""
    i = 0
    def add(txt, base=0):
        r = p.add_run(); r.text = txt
        r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color
        if base:
            r._r.get_or_add_rPr().set('baseline', str(base))
    buf = ""
    while i < len(latex):
        ch = latex[i]
        if ch in "_^":
            if buf: add(buf); buf = ""
            base = -25000 if ch == "_" else 30000
            i += 1
            if i < len(latex) and latex[i] == "{":
                j = latex.index("}", i); tok = latex[i+1:j]; i = j + 1
            else:
                tok = latex[i]; i += 1
            add(tok, base)
        else:
            buf += ch; i += 1
    if buf: add(buf)

def title(s, text):
    tb = s.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.9))
    p = tb.text_frame.paragraphs[0]; r = p.add_run(); r.text = text
    r.font.name = "Calibri"; r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = BLACK
    ln = s.shapes.add_shape(1, Inches(0.6), Inches(1.28), Inches(12.1), Pt(2))
    ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0xBB, 0xBB, 0xBB); ln.line.fill.background()

def line(s, top, latex, size=18, color=GREY, bold=False, font="Calibri", left=0.9, width=11.6, bullet=False):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.8))
    tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
    if bullet:
        rb = p.add_run(); rb.text = "•  "; rb.font.name = "Calibri"; rb.font.size = Pt(size)
        rb.font.color.rgb = color; rb.font.bold = bold
    _emit(p, latex, font, size, color, bold)
    return tb

def formula(s, top, latex, size=26, color=BLACK, left=1.4):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(10.8), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    _emit(p, latex, "Cambria Math", size, color, False)
    return tb

# ---------- 1. Title ----------
s = slide()
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.5), Inches(11.5), Inches(2.5))
p = tb.text_frame.paragraphs[0]; r = p.add_run()
r.text = "Relational Gradient Alignment (RGA)"
r.font.name = "Calibri"; r.font.size = Pt(38); r.font.bold = True; r.font.color.rgb = BLACK
p2 = tb.text_frame.add_paragraph(); r2 = p2.add_run()
r2.text = "Selecting KD data by pairwise alignment of two parameter spaces"
r2.font.name = "Calibri"; r2.font.size = Pt(22); r2.font.color.rgb = GREY
p3 = tb.text_frame.add_paragraph(); r3 = p3.add_run()
r3.text = "Teacher T (frozen)  →  Student S      signatures live in parameter space"
r3.font.name = "Calibri"; r3.font.size = Pt(16); r3.font.color.rgb = GREY

# ---------- 2. Setup & notation ----------
s = slide(); title(s, "Setup & notation")
line(s, 1.6, "Teacher T (frozen), student S; candidate pool of N samples. Select b samples for KD.", bullet=True)
line(s, 2.25, "p_T(x), p_S(x): teacher / student output distributions.", bullet=True)
line(s, 2.9, "θ_S, θ_T: student / teacher parameters (attention layers, all blocks).", bullet=True)
line(s, 3.55, "Challenge: θ_S and θ_T have different dimensions → cannot align the two parameter spaces directly.", bullet=True, color=RED, bold=True)
line(s, 4.35, "Idea: do pairwise / relative alignment — check if a group of samples follows the same relative", bullet=True, color=BLUE, bold=True)
line(s, 4.95, "pattern in the two parameter spaces.", left=1.25, color=BLUE, bold=True)

# ---------- 3. Step 1: signatures ----------
s = slide(); title(s, "Step 1 — two parameter-space signatures")
line(s, 1.55, "Student: gradient of the KD loss w.r.t. the student's deep parameters", bullet=True, color=BLACK, bold=True)
formula(s, 2.15, "g_S(x) = ∇_{θ_S} KL( p_T(x) ‖ p_S(x) )")
line(s, 3.05, "Teacher: gradient of the teacher's own loss w.r.t. the teacher's deep parameters (eNTK)", bullet=True, color=BLACK, bold=True)
formula(s, 3.65, "g_T(x) = ∇_{θ_T} ( − log p_T(y | x) )")
line(s, 4.55, "Random projection (Count-Sketch) to a common dimension d:", bullet=True)
formula(s, 5.1, "G_S = Π g_S ,   G_T = Π g_T      (both in parameter space)", size=22)
line(s, 6.0, "Fix vs. earlier version: teacher side used a last-layer feature (= last-layer eNTK, shallow);", bullet=True, color=RED)
line(s, 6.55, "now a deep parameter gradient → genuinely two parameter spaces.", left=1.25, color=RED)

# ---------- 4. Step 2: relational matrices ----------
s = slide(); title(s, "Step 2 — pairwise relation matrix in each space")
line(s, 1.6, "Pick m anchor samples a_1 … a_m. Build each space's Gram (pairwise) matrix:", bullet=True)
formula(s, 2.35, "K_S[i,j] = ⟨ G_S(x_i) , G_S(a_j) ⟩        (student parameter space)")
formula(s, 3.25, "K_T[i,j] = ⟨ G_T(x_i) , G_T(a_j) ⟩        (teacher parameter space)")
line(s, 4.25, "Double-center each (remove per-row / per-column bias, keep relative structure):", bullet=True)
formula(s, 5.0, "K′ = H K H ,     H = I − (1/m) · 1 1^T")
line(s, 6.0, "Each row K′[i,:] = how sample i relates to the anchors in that parameter space (its pattern).", bullet=True, color=GREY)

# ---------- 5. Step 3: residual ----------
s = slide(); title(s, "Step 3 — per-sample alignment residual")
line(s, 1.7, "For each sample, compare its pattern in the two parameter spaces:", bullet=True)
formula(s, 2.5, "r(i) = 1 − corr( K′_T[i,:] , K′_S[i,:] )", size=30)
line(s, 3.7, "r(i) small  →  same pattern in both parameter spaces  →  aligned  →  redundant (skip).", bullet=True, color=GREEN, bold=True)
line(s, 4.4, "r(i) large  →  different pattern  →  the teacher's parameter space encodes structure", bullet=True, color=RED, bold=True)
line(s, 5.0, "the student's does not yet  →  valuable (select).", left=1.25, color=RED, bold=True)
line(s, 5.9, "This is the CKA-style pairwise / relative alignment between the two parameter spaces.", bullet=True, color=BLUE, bold=True)

# ---------- 6. Step 4: selection ----------
s = slide(); title(s, "Step 4 — select via the alignment")
line(s, 1.7, "(optional) teacher-correctness gate — off for a clean teacher, on under a noised teacher.", bullet=True)
line(s, 2.5, "Take the top-3b samples by r; then D-optimal coverage on the anchor-space residual:", bullet=True)
formula(s, 3.3, "select b samples covering  R[i] = K′_T[i,:] − K′_S[i,:]")
line(s, 4.3, "Coverage keeps the selection diverse (not all the same kind of mis-alignment).", bullet=True, color=GREY)
line(s, 5.1, "Output: a coreset chosen purely by parameter-space alignment → distill the student on it.", bullet=True, color=BLACK, bold=True)

# ---------- 7. Pipeline at a glance ----------
s = slide(); title(s, "Whole pipeline at a glance")
formula(s, 1.8, "x  →  g_S(x)=∇_{θ_S} KL(p_T‖p_S) ,   g_T(x)=∇_{θ_T}(−log p_T)", size=20, left=0.8)
formula(s, 2.7, "→  project:  G_S , G_T", size=20, left=0.8)
formula(s, 3.6, "→  K_S = G_S G_S[A]^T ,   K_T = G_T G_T[A]^T   →  center  K′", size=20, left=0.8)
formula(s, 4.5, "→  r(i) = 1 − corr( K′_T[i,:] , K′_S[i,:] )", size=20, left=0.8)
formula(s, 5.4, "→  select high-r + coverage  →  b samples  →  KD-train S", size=20, left=0.8)

# ---------- 8. Why parameter-space alignment ----------
s = slide(); title(s, "Why this is parameter-space alignment")
rows = [("Requirement", "What the method does"),
        ("parameter space", "both g_S, g_T are deep parameter gradients"),
        ("different dims → pairwise/relative", "Gram K_S, K_T in each space; align the maps"),
        ("same pattern for a group", "r = 1 − corr of the two pattern rows"),
        ("use it to select data", "pick high-r (mis-aligned) samples + coverage")]
t = s.shapes.add_table(len(rows), 2, Inches(0.8), Inches(1.7), Inches(11.7), Inches(3.4)).table
t.columns[0].width = Inches(4.9); t.columns[1].width = Inches(6.8)
for ri, (a, b) in enumerate(rows):
    for ci, txt in enumerate((a, b)):
        c = t.cell(ri, ci); c.text = ""
        pp = c.text_frame.paragraphs[0]; _emit(pp, txt, "Calibri", 16 if ri else 17, BLACK if ri else RGBColor(0xFF,0xFF,0xFF), ri == 0)
        c.fill.solid(); c.fill.fore_color.rgb = BLUE if ri == 0 else (RGBColor(0xF2,0xF2,0xF2) if ri % 2 else RGBColor(0xFF,0xFF,0xFF))
line(s, 5.4, "Empirical: TAGCOS (same deep gradients, only clustered) and LESS (ICML'24, gradient-to-",
     bullet=True, color=GREEN, bold=True)
line(s, 6.0, "validation cosine) both lose to us at every budget → the teacher-relational alignment itself helps.",
     left=1.25, color=GREEN, bold=True)

out = "/data/tianhao/KD/RGA_method_for_advisor.pptx"
prs.save(out); print("wrote", out)
